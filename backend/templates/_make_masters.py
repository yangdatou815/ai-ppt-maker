"""Generate ``master.pptx`` for each template under ``backend/templates/``.

Each generated master is a **starter deck** users can open in PowerPoint /
Keynote and re-style. It carries:

* 16:9 slide size matching the renderer.
* The template's brand colours injected into the OOXML theme (so PowerPoint's
  *Design → Variants → Colors* shows the right palette).
* The template's heading / body fonts injected into the OOXML font scheme.
* Seven sample slides — one per layout the renderer emits — so the master
  doubles as a visual reference deck.

Run from the repository root::

    python backend/templates/_make_masters.py

Idempotent: re-running overwrites ``master.pptx`` in each template directory.

Limitations (intentional, deferred to M3-1 phase 2):
* We do not author *Slide Master / Slide Layouts* with bound placeholders;
  python-pptx cannot create custom layouts cleanly. The renderer remains
  programmatic for v0.5.x; this file only provides a brand-tinted starter.
"""

from __future__ import annotations

import re
import sys
from pathlib import Path

import yaml
from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).resolve().parent
SLIDE_W = Emu(12192000)  # 13.333"
SLIDE_H = Emu(6858000)  #  7.5"  (16:9)

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
NSMAP = {"a": A_NS}

_HEX_RE = re.compile(r"#?([0-9a-fA-F]{6})")


def _hex(value: str | None, fallback: str) -> str:
    if not value:
        return fallback
    m = _HEX_RE.search(value)
    return (m.group(1) if m else fallback).upper()


def _rgb(hex6: str) -> RGBColor:
    return RGBColor(int(hex6[0:2], 16), int(hex6[2:4], 16), int(hex6[4:6], 16))


def _patch_theme(prs: Presentation, theme: dict, fonts: dict) -> None:
    """Rewrite <a:clrScheme> and <a:fontScheme> inside ppt/theme/theme1.xml.

    Maps our token vocabulary onto OOXML's fixed slot names so PowerPoint's
    Design picker reflects the brand identity:

        primary    -> accent1
        accent     -> accent2
        primary-2  -> dk2
        text       -> dk1
        text-mute  -> accent3
        background -> lt1
        neutral    -> lt2
    """
    primary = _hex(theme.get("primary"), "1F6FEB")
    accent = _hex(theme.get("accent"), "39D0D8")
    primary2 = _hex(theme.get("primary-2"), "0D1117")
    text = _hex(theme.get("text"), "0B0B0F")
    text_mute = _hex(theme.get("text-mute"), "6B7280")
    bg = _hex(theme.get("background") or theme.get("neutral"), "FFFFFF")
    neutral = _hex(theme.get("neutral") or theme.get("background"), "F4F1EA")
    heading_font = fonts.get("heading", "Inter")
    body_font = fonts.get("body", "Inter")

    prs.slide_masters[0].element.getroottree().getroot()
    # Walk to the actual theme XML part via the package.
    for part in prs.part.package.iter_parts():
        if part.partname.endswith("theme1.xml"):
            theme_xml = etree.fromstring(part.blob)
            break
    else:
        return

    # --- Color scheme ---
    clr = theme_xml.find(".//a:clrScheme", NSMAP)
    if clr is not None:
        slot_map = {
            "dk1": text,
            "lt1": bg,
            "dk2": primary2,
            "lt2": neutral,
            "accent1": primary,
            "accent2": accent,
            "accent3": text_mute,
        }
        for child in list(clr):
            tag = etree.QName(child).localname
            if tag in slot_map:
                # Clear existing child colour element, set srgbClr value.
                for sub in list(child):
                    child.remove(sub)
                srgb = etree.SubElement(child, f"{{{A_NS}}}srgbClr")
                srgb.set("val", slot_map[tag])

    # --- Font scheme ---
    font_scheme = theme_xml.find(".//a:fontScheme", NSMAP)
    if font_scheme is not None:
        for major_minor in ("majorFont", "minorFont"):
            fnode = font_scheme.find(f"a:{major_minor}", NSMAP)
            if fnode is None:
                continue
            font_name = heading_font if major_minor == "majorFont" else body_font
            latin = fnode.find("a:latin", NSMAP)
            if latin is not None:
                latin.set("typeface", font_name)
            ea = fnode.find("a:ea", NSMAP)
            if ea is not None:
                ea.set("typeface", font_name)

    # Write back.
    new_blob = etree.tostring(theme_xml, xml_declaration=True, encoding="UTF-8", standalone=True)
    part._blob = new_blob


def _add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    font_name="Inter",
    size_pt=18,
    bold=False,
    color="0B0B0F",
    align="left",
):
    from pptx.enum.text import PP_ALIGN

    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    return box


def _accent_rule(slide, left, top, width, color):
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Emu(45720))  # 0.05"
    rule.line.fill.background()
    rule.fill.solid()
    rule.fill.fore_color.rgb = _rgb(color)
    return rule


def _bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color)


def _add_sample_slides(prs: Presentation, theme: dict, fonts: dict, name: str) -> None:
    primary = _hex(theme.get("primary"), "1F6FEB")
    primary2 = _hex(theme.get("primary-2"), "0D1117")
    accent = _hex(theme.get("accent"), "39D0D8")
    text = _hex(theme.get("text"), "0B0B0F")
    text_mute = _hex(theme.get("text-mute"), "6B7280")
    bg = _hex(theme.get("background") or theme.get("neutral"), "FFFFFF")
    heading_font = fonts.get("heading", "Inter")
    body_font = fonts.get("body", "Inter")
    blank = prs.slide_layouts[6]

    # 1. Cover -----------------------------------------------------------
    s = prs.slides.add_slide(blank)
    _bg(s, primary2 if name == "executive-dark" else bg)
    title_color = bg if name == "executive-dark" else text
    sub_color = text_mute
    _accent_rule(s, Inches(0.6), Inches(2.2), Inches(0.8), accent)
    _add_text(
        s,
        Inches(0.6),
        Inches(2.4),
        Inches(11),
        Inches(1.5),
        "Cover Slide Title",
        font_name=heading_font,
        size_pt=48,
        bold=True,
        color=title_color,
    )
    _add_text(
        s,
        Inches(0.6),
        Inches(4.0),
        Inches(11),
        Inches(0.6),
        "Subtitle / one-line value proposition",
        font_name=body_font,
        size_pt=20,
        color=sub_color,
    )
    _add_text(
        s,
        Inches(0.6),
        Inches(6.7),
        Inches(11),
        Inches(0.4),
        f"{name} · master.pptx",
        font_name=body_font,
        size_pt=11,
        color=sub_color,
    )

    # 2. TOC -------------------------------------------------------------
    s = prs.slides.add_slide(blank)
    _bg(s, bg)
    _add_text(
        s,
        Inches(0.6),
        Inches(0.5),
        Inches(11),
        Inches(0.8),
        "Agenda",
        font_name=heading_font,
        size_pt=36,
        bold=True,
        color=text,
    )
    _accent_rule(s, Inches(0.6), Inches(1.35), Inches(0.8), primary)
    items = ["1. Background", "2. Approach", "3. Results", "4. Next steps"]
    for i, t in enumerate(items):
        _add_text(
            s,
            Inches(0.8),
            Inches(1.9 + i * 0.7),
            Inches(10),
            Inches(0.6),
            t,
            font_name=body_font,
            size_pt=22,
            color=text,
        )

    # 3. Section divider -------------------------------------------------
    s = prs.slides.add_slide(blank)
    _bg(s, primary)
    _add_text(
        s,
        Inches(0.6),
        Inches(2.5),
        Inches(2),
        Inches(2),
        "01",
        font_name=heading_font,
        size_pt=120,
        bold=True,
        color=accent,
    )
    _add_text(
        s,
        Inches(2.6),
        Inches(3.2),
        Inches(10),
        Inches(1.2),
        "Section Divider",
        font_name=heading_font,
        size_pt=44,
        bold=True,
        color=bg,
    )

    # 4. content-bullets -------------------------------------------------
    s = prs.slides.add_slide(blank)
    _bg(s, bg)
    _add_text(
        s,
        Inches(0.6),
        Inches(0.5),
        Inches(11),
        Inches(0.8),
        "Bullets layout",
        font_name=heading_font,
        size_pt=32,
        bold=True,
        color=text,
    )
    _accent_rule(s, Inches(0.6), Inches(1.3), Inches(0.6), primary)
    bullets = [
        "First main point with a short clarifying clause",
        "Second point — concise, action-oriented",
        "Third point references prior context",
        "Fourth point closes the argument",
    ]
    for i, b in enumerate(bullets):
        _add_text(
            s,
            Inches(0.8),
            Inches(1.8 + i * 0.7),
            Inches(11),
            Inches(0.6),
            f"• {b}",
            font_name=body_font,
            size_pt=20,
            color=text,
        )

    # 5. content-image ---------------------------------------------------
    s = prs.slides.add_slide(blank)
    _bg(s, bg)
    _add_text(
        s,
        Inches(0.6),
        Inches(0.5),
        Inches(11),
        Inches(0.8),
        "Image layout",
        font_name=heading_font,
        size_pt=32,
        bold=True,
        color=text,
    )
    _accent_rule(s, Inches(0.6), Inches(1.3), Inches(0.6), primary)
    for i, b in enumerate(
        [
            "• Bullet on the left half",
            "• Image on the right half",
            "• Aspect ratio preserved",
            "• Golden-ratio framing",
        ]
    ):
        _add_text(
            s,
            Inches(0.8),
            Inches(1.9 + i * 0.7),
            Inches(5.4),
            Inches(0.6),
            b,
            font_name=body_font,
            size_pt=18,
            color=text,
        )
    placeholder = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7.0),
        Inches(1.9),
        Inches(5.9),
        Inches(4.0),
    )
    placeholder.fill.solid()
    placeholder.fill.fore_color.rgb = _rgb(text_mute)
    placeholder.line.fill.background()
    _add_text(
        s,
        Inches(7.0),
        Inches(3.7),
        Inches(5.9),
        Inches(0.5),
        "[ image placeholder ]",
        font_name=body_font,
        size_pt=14,
        color=bg,
        align="center",
    )

    # 6. content-table ---------------------------------------------------
    s = prs.slides.add_slide(blank)
    _bg(s, bg)
    _add_text(
        s,
        Inches(0.6),
        Inches(0.5),
        Inches(11),
        Inches(0.8),
        "Table layout",
        font_name=heading_font,
        size_pt=32,
        bold=True,
        color=text,
    )
    _accent_rule(s, Inches(0.6), Inches(1.3), Inches(0.6), primary)
    rows, cols = 4, 3
    table_shape = s.shapes.add_table(
        rows, cols, Inches(0.8), Inches(1.9), Inches(11.6), Inches(3.5)
    )
    table = table_shape.table
    headers = ["Metric", "Q3", "Q4"]
    data = [["Revenue", "$1.2M", "$1.8M"], ["Users", "12k", "21k"], ["NPS", "42", "57"]]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(primary)
        cell.text = ""
        run = cell.text_frame.paragraphs[0].add_run()
        run.text = h
        run.font.bold = True
        run.font.size = Pt(16)
        run.font.color.rgb = _rgb(bg)
        run.font.name = heading_font
    for r, row in enumerate(data, start=1):
        for c, v in enumerate(row):
            cell = table.cell(r, c)
            cell.text = ""
            run = cell.text_frame.paragraphs[0].add_run()
            run.text = v
            run.font.size = Pt(14)
            run.font.name = body_font
            run.font.color.rgb = _rgb(text)

    # 7. Closing ---------------------------------------------------------
    s = prs.slides.add_slide(blank)
    _bg(s, primary2 if name == "executive-dark" else bg)
    title_color = bg if name == "executive-dark" else text
    _accent_rule(s, Inches(0.6), Inches(3.0), Inches(0.8), accent)
    _add_text(
        s,
        Inches(0.6),
        Inches(3.2),
        Inches(11),
        Inches(1.5),
        "Thank you",
        font_name=heading_font,
        size_pt=64,
        bold=True,
        color=title_color,
    )
    _add_text(
        s,
        Inches(0.6),
        Inches(4.6),
        Inches(11),
        Inches(0.6),
        "Questions & discussion",
        font_name=body_font,
        size_pt=22,
        color=text_mute,
    )


def build_master(template_dir: Path) -> Path:
    mapping_file = template_dir / "layout-mapping.yaml"
    data = yaml.safe_load(mapping_file.read_text(encoding="utf-8")) or {}
    theme = data.get("theme", {}) or {}
    fonts = data.get("fonts", {}) or {}

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _patch_theme(prs, theme, fonts)
    _add_sample_slides(prs, theme, fonts, template_dir.name)

    out = template_dir / "master.pptx"
    prs.save(str(out))
    return out


def main() -> int:
    targets = [
        d for d in sorted(HERE.iterdir()) if d.is_dir() and (d / "layout-mapping.yaml").is_file()
    ]
    if not targets:
        print(f"[!] no template directories under {HERE}", file=sys.stderr)
        return 1
    for d in targets:
        out = build_master(d)
        size_kb = out.stat().st_size // 1024
        print(
            f"[OK] {out.relative_to(HERE.parent.parent)}  ({size_kb} KB, "
            f"{len(Presentation(str(out)).slides)} slides)"
        )
    return 0


if __name__ == "__main__":
    sys.exit(main())
