"""Unit tests for app.render.pptx_renderer."""
from __future__ import annotations

import io
import struct
import zlib
from pathlib import Path

import pptx
import pytest

from app.render.pptx_renderer import (
    _CONTENT_HEIGHT,
    _CONTENT_WIDTH,
    _GUTTER,
    _IMG_MAX_FRAC,
    _image_box_for_aspect,
    _pick_layout,
    _table_column_widths,
    _Theme,
    render_outline,
)
from app.schemas.outline import (
    Bullet,
    ImageRef,
    OutlineDoc,
    Section,
    TableData,
)
from app.schemas.template import TemplateInfo


def _doc(n_sections: int = 2) -> OutlineDoc:
    return OutlineDoc(
        title="Demo Title",
        subtitle="Demo Subtitle",
        language="zh",
        sections=[
            Section(
                heading=f"Section {i + 1}",
                bullets=[Bullet(text=f"point {i + 1}", emphasis=(i == 0))],
                speaker_notes=f"note {i + 1}" if i == 0 else None,
            )
            for i in range(n_sections)
        ],
        cover_meta={"company": "ACME", "author": "YDT"},
    )


def _template(theme: dict[str, str] | None = None) -> TemplateInfo:
    return TemplateInfo(
        name="t",
        display_name="t",
        description="",
        tags=[],
        theme=theme
        or {
            "primary": "#0B1F3A",
            "primary-2": "#050B14",
            "accent": "#C9A24E",
            "text": "#F4F1EA",
            "text-mute": "#A3A6AD",
            "neutral": "#F4F1EA",
        },
        fonts={"heading": "Calibri", "body": "Calibri"},
        has_master=False,
        thumbnail_url=None,
    )


def test_render_outline_returns_valid_pptx_bytes():
    data = render_outline(_doc(3), _template())
    assert data[:4] == b"PK\x03\x04"
    prs = pptx.Presentation(io.BytesIO(data))
    # cover + N sections + closing
    assert len(prs.slides) == 3 + 2


def test_render_outline_speaker_notes_go_to_notes_pane():
    data = render_outline(_doc(1), _template())
    prs = pptx.Presentation(io.BytesIO(data))
    # First content slide is index 1 (after cover)
    notes = prs.slides[1].notes_slide.notes_text_frame.text
    assert "note 1" in notes


def test_theme_falls_back_when_hex_garbage():
    bad = _template(theme={"primary": "not-a-color", "accent": "ZZZZZZ"})
    # Must not raise — defaults kick in.
    theme = _Theme.from_template(bad)
    assert theme.primary is not None
    assert theme.accent is not None


def test_render_outline_handles_empty_sections():
    doc = OutlineDoc(title="Empty", language="en", sections=[])
    data = render_outline(doc, _template())
    prs = pptx.Presentation(io.BytesIO(data))
    # cover + closing only
    assert len(prs.slides) == 2


# ---------------------------------------------------------------------------
# M2-3: layout_hint dispatch + image / table rendering
# ---------------------------------------------------------------------------


def _png_1x1() -> bytes:
    """Minimal valid 1x1 PNG. Lets us avoid pulling Pillow into test deps."""
    return _png(1, 1)


def _png(width: int, height: int) -> bytes:
    """Minimal valid solid-colour PNG of the given dimensions."""
    sig = b"\x89PNG\r\n\x1a\n"

    def _chunk(typ: bytes, data: bytes) -> bytes:
        return (
            struct.pack(">I", len(data))
            + typ
            + data
            + struct.pack(">I", zlib.crc32(typ + data) & 0xFFFFFFFF)
        )

    ihdr = _chunk(b"IHDR", struct.pack(">IIBBBBB", width, height, 8, 2, 0, 0, 0))
    # filter byte (0) per row + RGB pixels
    row = b"\x00" + b"\xff\x00\x00" * width
    raw = row * height
    idat = _chunk(b"IDAT", zlib.compress(raw))
    iend = _chunk(b"IEND", b"")
    return sig + ihdr + idat + iend


def test_pick_layout_explicit_hint_wins():
    s = Section(heading="h", layout_hint="content-image",
                image=ImageRef(file_id="x.png"))
    assert _pick_layout(s) == "content-image"


def test_pick_layout_falls_back_when_hint_mismatches_payload():
    # hint says image but no image attached → fall through to inference (bullets)
    s = Section(heading="h", layout_hint="content-image",
                bullets=[Bullet(text="b")])
    assert _pick_layout(s) == "content-bullets"


def test_pick_layout_infers_from_payload():
    assert _pick_layout(Section(heading="h", image=ImageRef(file_id="a.png"))) == "content-image"
    assert _pick_layout(
        Section(heading="h", table=TableData(headers=["a"], rows=[["1"]]))
    ) == "content-table"
    assert _pick_layout(Section(heading="h")) == "content-bullets"


def test_render_table_section_creates_native_pptx_table():
    doc = OutlineDoc(
        title="T",
        language="en",
        sections=[
            Section(
                heading="Numbers",
                table=TableData(
                    headers=["Q", "Rev"],
                    rows=[["Q1", "10"], ["Q2", "20"], ["Q3", "30"]],
                    caption="quarterly",
                ),
            )
        ],
    )
    data = render_outline(doc, _template())
    prs = pptx.Presentation(io.BytesIO(data))
    table_slide = prs.slides[1]
    tables = [s for s in table_slide.shapes if s.has_table]
    assert len(tables) == 1
    tbl = tables[0].table
    # header row + 3 data rows
    assert len(tbl.rows) == 4
    assert len(tbl.columns) == 2
    assert tbl.cell(0, 0).text_frame.text == "Q"
    assert tbl.cell(2, 1).text_frame.text == "20"


def test_render_image_section_with_real_file_embeds_picture(tmp_path: Path):
    uploads = tmp_path / "uploads"
    uploads.mkdir()
    (uploads / "pic.png").write_bytes(_png_1x1())

    doc = OutlineDoc(
        title="I",
        language="en",
        sections=[
            Section(
                heading="See chart",
                bullets=[Bullet(text="left side")],
                image=ImageRef(file_id="pic.png", caption="cap"),
            )
        ],
    )
    data = render_outline(doc, _template(), uploads_dir=uploads)
    prs = pptx.Presentation(io.BytesIO(data))
    pics = [s for s in prs.slides[1].shapes if s.shape_type == 13]  # MSO_SHAPE_TYPE.PICTURE
    assert len(pics) == 1


def test_render_image_section_placeholder_when_file_missing(tmp_path: Path):
    uploads = tmp_path / "uploads"
    uploads.mkdir()
    doc = OutlineDoc(
        title="I",
        language="en",
        sections=[
            Section(heading="x", image=ImageRef(file_id="ghost.png")),
        ],
    )
    data = render_outline(doc, _template(), uploads_dir=uploads)
    prs = pptx.Presentation(io.BytesIO(data))
    # No PICTURE shapes — only the placeholder rect + text.
    pics = [s for s in prs.slides[1].shapes if s.shape_type == 13]
    assert pics == []
    # And the slide must still render (cover + 1 + closing).
    assert len(prs.slides) == 3


def test_render_image_section_rejects_path_traversal(tmp_path: Path):
    uploads = tmp_path / "uploads"
    uploads.mkdir()
    secret = tmp_path / "secret.png"
    secret.write_bytes(_png_1x1())

    doc = OutlineDoc(
        title="I",
        language="en",
        sections=[
            Section(heading="x", image=ImageRef(file_id="../secret.png")),
        ],
    )
    data = render_outline(doc, _template(), uploads_dir=uploads)
    prs = pptx.Presentation(io.BytesIO(data))
    pics = [s for s in prs.slides[1].shapes if s.shape_type == 13]
    assert pics == []  # traversal rejected


def test_render_table_with_truncation_warning():
    rows = [[f"r{i}", str(i)] for i in range(20)]
    doc = OutlineDoc(
        title="T",
        language="en",
        sections=[
            Section(
                heading="Big",
                table=TableData(headers=["k", "v"], rows=rows),
            )
        ],
    )
    data = render_outline(doc, _template())
    prs = pptx.Presentation(io.BytesIO(data))
    tbls = [s for s in prs.slides[1].shapes if s.has_table]
    assert len(tbls) == 1
    # Header + 14 cap = 15 rows, regardless of input being 20.
    assert len(tbls[0].table.rows) == 15


# ---------------------------------------------------------------------------
# Image aspect-ratio + golden-ratio layout (M2-3 follow-up)
# ---------------------------------------------------------------------------


def test_image_box_for_aspect_landscape_clamped_to_golden_max():
    # Very wide panorama (4:1) — width would exceed φ-bound, must clamp.
    text_w, img_w, img_h, _ = _image_box_for_aspect(4.0)
    inner_w = _CONTENT_WIDTH - _GUTTER
    assert img_w == pytest.approx(inner_w * _IMG_MAX_FRAC, rel=1e-6)
    # Aspect preserved: img_w / img_h ≈ 4.0
    assert img_w / img_h == pytest.approx(4.0, rel=1e-6)
    # Text column gets the remainder.
    assert text_w == pytest.approx(_CONTENT_WIDTH - _GUTTER - img_w, rel=1e-6)


def test_image_box_for_aspect_portrait_clamped_to_golden_min():
    # Tall portrait (1:3) — width would be below φ-min. Even after clamping up,
    # the resulting height would exceed the band, so the box ends up height-
    # limited (img_h == content height) with proportionally narrow width.
    # Aspect must still be preserved exactly.
    text_w, img_w, img_h, top_off = _image_box_for_aspect(1.0 / 3.0)
    assert img_h == pytest.approx(_CONTENT_HEIGHT, rel=1e-6)
    assert img_w / img_h == pytest.approx(1.0 / 3.0, rel=1e-6)
    assert top_off == pytest.approx(0.0, abs=1e-6)
    assert text_w > 0


def test_image_box_for_aspect_moderate_portrait_clamps_width():
    # Portrait 2:3 — fit-to-height yields img_w = 4.55 * 2/3 ≈ 3.03,
    # which IS below min. Clamping recomputes height = min_w * 3/2 ≈ 6.6,
    # which still exceeds band. So we again fall back to height-limited.
    # 9:16 phone aspect (≈0.5625) is a more realistic clamp scenario:
    # fit-height gives img_w ≈ 2.56 < min ≈ 4.40 → clamp width up to min,
    # img_h = 4.40 / 0.5625 ≈ 7.82 > band → height-limited again.
    # Truly: anything below ~0.97 aspect ends up height-limited.
    text_w, img_w, img_h, _ = _image_box_for_aspect(0.5625)
    assert img_h == pytest.approx(_CONTENT_HEIGHT, rel=1e-6)
    assert img_w / img_h == pytest.approx(0.5625, rel=1e-6)
    assert text_w > 0


def test_image_box_for_aspect_square_fits_height_and_centres_columns():
    text_w, img_w, img_h, top_off = _image_box_for_aspect(1.0)
    # Square within bounds: limited by height, so img_h == content height.
    assert img_h == pytest.approx(_CONTENT_HEIGHT, rel=1e-6)
    assert img_w == pytest.approx(_CONTENT_HEIGHT, rel=1e-6)
    assert top_off == pytest.approx(0.0, abs=1e-6)
    assert text_w > 0


def test_render_image_section_preserves_aspect_ratio_landscape(tmp_path: Path):
    uploads = tmp_path / "uploads"
    uploads.mkdir()
    (uploads / "wide.png").write_bytes(_png(200, 100))  # 2:1
    doc = OutlineDoc(
        title="I", language="en",
        sections=[Section(heading="x", image=ImageRef(file_id="wide.png"))],
    )
    data = render_outline(doc, _template(), uploads_dir=uploads)
    prs = pptx.Presentation(io.BytesIO(data))
    pics = [s for s in prs.slides[1].shapes if s.shape_type == 13]
    assert len(pics) == 1
    pic = pics[0]
    # Source aspect 2.0 must be preserved on the rendered shape.
    assert pic.width / pic.height == pytest.approx(2.0, rel=1e-3)


def test_render_image_section_preserves_aspect_ratio_portrait(tmp_path: Path):
    uploads = tmp_path / "uploads"
    uploads.mkdir()
    (uploads / "tall.png").write_bytes(_png(100, 200))  # 1:2
    doc = OutlineDoc(
        title="I", language="en",
        sections=[Section(heading="x", image=ImageRef(file_id="tall.png"))],
    )
    data = render_outline(doc, _template(), uploads_dir=uploads)
    prs = pptx.Presentation(io.BytesIO(data))
    pics = [s for s in prs.slides[1].shapes if s.shape_type == 13]
    assert len(pics) == 1
    pic = pics[0]
    assert pic.width / pic.height == pytest.approx(0.5, rel=1e-3)


# ---------------------------------------------------------------------------
# Table content-adaptive width
# ---------------------------------------------------------------------------


def test_table_column_widths_short_content_stays_compact():
    widths = _table_column_widths(["a", "b"], [["1", "2"], ["3", "4"]])
    # Two short columns should be near the per-column minimum (0.9").
    assert all(w == pytest.approx(0.9, rel=1e-6) for w in widths)
    assert sum(widths) < 3.0


def test_table_column_widths_long_content_expands_proportionally():
    widths = _table_column_widths(
        ["k", "description"],
        [["x", "a long descriptive sentence about something"]],
    )
    # Second column should be noticeably wider than the first.
    assert widths[1] > widths[0] * 2


def test_table_column_widths_scaled_down_when_total_exceeds_max():
    # 5 cols × ~5" each would be 25", but max_total is 11.9 → must scale down.
    widths = _table_column_widths(
        ["a" * 60] * 5,
        [["x" * 60] * 5],
        max_total_in=11.9,
    )
    assert sum(widths) == pytest.approx(11.9, rel=1e-6)


def test_render_table_section_centres_compact_table_horizontally():
    doc = OutlineDoc(
        title="T", language="en",
        sections=[
            Section(
                heading="Tiny",
                table=TableData(headers=["a", "b"], rows=[["1", "2"]]),
            )
        ],
    )
    data = render_outline(doc, _template(), uploads_dir=None)
    prs = pptx.Presentation(io.BytesIO(data))
    tables = [s for s in prs.slides[1].shapes if s.has_table]
    assert len(tables) == 1
    shape = tables[0]
    # Compact table should be centred (left margin > 0.7 default).
    left_in = shape.left / 914400  # EMUs per inch
    width_in = shape.width / 914400
    # Centred means left + width/2 ≈ slide_w / 2
    assert (left_in + width_in / 2) == pytest.approx(13.333 / 2, abs=0.05)


# --- M3-1 phase 2: master.pptx loading -------------------------------------

_MASTERS_ROOT = Path(__file__).resolve().parents[2] / "templates"


def _expected_slide_count(n_sections: int) -> int:
    # cover + N section slides + closing
    return 1 + n_sections + 1


def test_render_with_master_strips_sample_slides_and_keeps_user_slides():
    master = _MASTERS_ROOT / "tech-blue" / "master.pptx"
    assert master.is_file(), f"missing fixture: {master}"
    n = 3
    data = render_outline(_doc(n), _template(), master_path=master)
    prs = pptx.Presentation(io.BytesIO(data))
    # The master ships 7 sample slides; output must contain only the user's.
    assert len(prs.slides) == _expected_slide_count(n)


def test_render_without_master_path_matches_blank_baseline():
    # Default behaviour (no master_path) must remain identical to before.
    data = render_outline(_doc(2), _template())
    prs = pptx.Presentation(io.BytesIO(data))
    assert len(prs.slides) == _expected_slide_count(2)


def test_render_with_missing_master_path_falls_back_silently(tmp_path):
    bogus = tmp_path / "does-not-exist.pptx"
    data = render_outline(_doc(1), _template(), master_path=bogus)
    prs = pptx.Presentation(io.BytesIO(data))
    assert len(prs.slides) == _expected_slide_count(1)


def test_render_with_corrupt_master_falls_back_silently(tmp_path):
    bad = tmp_path / "corrupt.pptx"
    bad.write_bytes(b"this is not a valid pptx zip")
    data = render_outline(_doc(1), _template(), master_path=bad)
    prs = pptx.Presentation(io.BytesIO(data))
    assert len(prs.slides) == _expected_slide_count(1)

