"""Programmatic PPTX renderer (M2-2 / M2-3).

Renders an :class:`OutlineDoc` to a python-pptx :class:`Presentation` using the
template's theme tokens (no real ``master.pptx`` required for v0.3.0). Each
section becomes one of:

- ``cover``: title + subtitle + accent rule
- ``content-bullets``: numbered heading + bullet column with emphasis-bold
- ``content-image``: heading + bullets on the left, embedded image on the
  right (or a placeholder rect when the file is missing)
- ``content-table``: heading + native pptx table with header-row styling
- ``closing``: thank-you slide with the same accent treatment as cover

Layout is chosen by :func:`_pick_layout`: explicit ``section.layout_hint``
wins; otherwise we infer from the presence of ``image`` / ``table``.

Design tokens come from ``backend/templates/<name>/layout-mapping.yaml``
(``theme:`` + ``fonts:``). Hex colours are parsed defensively — unknown keys
fall back to a sensible default rather than crashing the request.
"""
from __future__ import annotations

import io
import logging
import re
from dataclasses import dataclass
from pathlib import Path

from PIL import Image, UnidentifiedImageError
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from app.schemas.outline import OutlineDoc, Section
from app.schemas.template import TemplateInfo

log = logging.getLogger(__name__)

_HEX_RE = re.compile(r"#?([0-9a-fA-F]{6})")
_DEFAULT_THEME = {
    "primary": "0B1F3A",
    "primary-2": "050B14",
    "accent": "C9A24E",
    "text": "F4F1EA",
    "text-mute": "A3A6AD",
    "neutral": "F4F1EA",
}
_DEFAULT_FONTS = {
    "heading": "Calibri",
    "body": "Calibri",
}


@dataclass(frozen=True)
class _Theme:
    primary: RGBColor
    primary_2: RGBColor
    accent: RGBColor
    text: RGBColor
    text_mute: RGBColor
    neutral: RGBColor
    heading_font: str
    body_font: str

    @classmethod
    def from_template(cls, t: TemplateInfo) -> _Theme:
        def _rgb(key: str) -> RGBColor:
            raw = t.theme.get(key) or _DEFAULT_THEME[key]
            m = _HEX_RE.search(raw)
            return RGBColor.from_string((m.group(1) if m else _DEFAULT_THEME[key]).upper())

        return cls(
            primary=_rgb("primary"),
            primary_2=_rgb("primary-2"),
            accent=_rgb("accent"),
            text=_rgb("text"),
            text_mute=_rgb("text-mute"),
            neutral=_rgb("neutral"),
            heading_font=t.fonts.get("heading", _DEFAULT_FONTS["heading"]),
            body_font=t.fonts.get("body", _DEFAULT_FONTS["body"]),
        )


# 16:9, 13.333" x 7.5"
_SLIDE_W = Inches(13.333)
_SLIDE_H = Inches(7.5)


def _set_slide_bg(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _add_rect(slide, left, top, width, height, rgb: RGBColor):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    return shape


def _add_text(
    slide,
    left,
    top,
    width,
    height,
    *,
    text: str,
    font: str,
    size_pt: int,
    color: RGBColor,
    bold: bool = False,
    align: PP_ALIGN | None = None,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = font
    f.size = Pt(size_pt)
    f.bold = bold
    f.color.rgb = color
    return tb


def _cover_slide(prs: Presentation, theme: _Theme, doc: OutlineDoc) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_slide_bg(slide, theme.primary)

    # Accent rule above the title
    _add_rect(slide, Inches(1.0), Inches(2.4), Inches(0.6), Inches(0.06), theme.accent)
    _add_text(
        slide, Inches(1.0), Inches(2.55), Inches(11.3), Inches(1.6),
        text=doc.title,
        font=theme.heading_font, size_pt=54, color=theme.text, bold=True,
    )
    if doc.subtitle:
        _add_text(
            slide, Inches(1.0), Inches(4.2), Inches(11.3), Inches(0.8),
            text=doc.subtitle,
            font=theme.body_font, size_pt=22, color=theme.text_mute,
        )
    # Footer line
    parts = [
        doc.cover_meta.get("company"),
        doc.cover_meta.get("author"),
        doc.cover_meta.get("date"),
    ]
    footer = "  ·  ".join(p for p in parts if p)
    if footer:
        _add_text(
            slide, Inches(1.0), Inches(6.6), Inches(11.3), Inches(0.4),
            text=footer,
            font=theme.body_font, size_pt=12, color=theme.text_mute,
        )


def _section_header(slide, theme: _Theme, idx: int, total: int, heading: str) -> None:
    """Common top-of-slide header: number gutter + heading + accent rule."""
    _add_text(
        slide, Inches(0.7), Inches(0.55), Inches(2.0), Inches(0.5),
        text=f"{idx:02d} / {total:02d}",
        font=theme.body_font, size_pt=12, color=theme.accent, bold=True,
    )
    _add_text(
        slide, Inches(0.7), Inches(1.05), Inches(11.9), Inches(1.0),
        text=heading,
        font=theme.heading_font, size_pt=34, color=theme.primary, bold=True,
    )
    _add_rect(slide, Inches(0.7), Inches(2.05), Inches(0.6), Inches(0.05), theme.accent)


def _draw_bullets(slide, theme: _Theme, section: Section, *, left, top, width, height) -> None:
    if not section.bullets:
        return
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(section.bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = "•  " + b.text
        run.font.name = theme.body_font
        run.font.size = Pt(20)
        run.font.bold = bool(b.emphasis)
        run.font.color.rgb = theme.primary if b.emphasis else theme.text_mute
        if b.note:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.LEFT
            r2 = p2.add_run()
            r2.text = "    — " + b.note
            r2.font.name = theme.body_font
            r2.font.size = Pt(14)
            r2.font.italic = True
            r2.font.color.rgb = theme.text_mute


def _bullets_section_slide(slide, theme: _Theme, section: Section) -> None:
    _draw_bullets(
        slide, theme, section,
        left=Inches(0.9), top=Inches(2.45), width=Inches(11.5), height=Inches(4.5),
    )


def _resolve_upload(file_id: str | None, uploads_dir: Path | None) -> Path | None:
    """Map a section ``file_id`` to a real upload, with traversal guard."""
    if not file_id or uploads_dir is None:
        return None
    candidate = (uploads_dir / file_id).resolve()
    try:
        candidate.relative_to(uploads_dir.resolve())
    except ValueError:
        log.warning("rejecting upload path outside uploads_dir: %s", file_id)
        return None
    return candidate if candidate.is_file() else None


# Golden-ratio constants for image-slide layout. The available content band
# below the slide header is ~4.55" tall and ~11.93" wide. We allow the image
# box to occupy between 38.2% and 61.8% of that width, so the text/image
# proportion always falls inside the golden window even for extreme aspects.
_PHI = 1.618
_CONTENT_LEFT = 0.7      # inches
_CONTENT_TOP = 2.45      # inches
_CONTENT_WIDTH = 11.93   # inches
_CONTENT_HEIGHT = 4.55   # inches
_GUTTER = 0.4            # inches between text column and image
_IMG_MIN_FRAC = 1.0 / (1.0 + _PHI)   # ≈ 0.382
_IMG_MAX_FRAC = _PHI / (1.0 + _PHI)  # ≈ 0.618


def _image_box_for_aspect(aspect: float) -> tuple[float, float, float, float]:
    """Return ``(text_w, img_w, img_h, img_top_offset)`` in inches for a given image aspect.

    Image is fit-into a box that preserves ``aspect`` (= width/height). Box width
    is clamped to the golden-ratio bounds so the page composition stays balanced.
    Image is vertically centred within the content band.
    """
    inner_w = _CONTENT_WIDTH - _GUTTER
    # Try fit-to-height first.
    img_h = _CONTENT_HEIGHT
    img_w = img_h * aspect
    img_w_min = inner_w * _IMG_MIN_FRAC
    img_w_max = inner_w * _IMG_MAX_FRAC
    if img_w > img_w_max:
        img_w = img_w_max
        img_h = img_w / aspect
    elif img_w < img_w_min:
        img_w = img_w_min
        img_h = img_w / aspect
        if img_h > _CONTENT_HEIGHT:
            # Very tall portrait — recompute from height again with new width cap.
            img_h = _CONTENT_HEIGHT
            img_w = img_h * aspect
    text_w = _CONTENT_WIDTH - _GUTTER - img_w
    img_top_offset = max(0.0, (_CONTENT_HEIGHT - img_h) / 2.0)
    return text_w, img_w, img_h, img_top_offset


def _image_section_slide(
    slide, theme: _Theme, section: Section, uploads_dir: Path | None,
) -> None:
    """Render an image section.

    Images preserve their original aspect ratio (no stretching). The text
    column width adapts to the image so the overall composition approximates
    the golden ratio (φ ≈ 1.618).
    """
    file_id = section.image.file_id if section.image else None
    img_path = _resolve_upload(file_id, uploads_dir)

    # Probe real dimensions to drive the layout. If the file is missing or
    # unreadable we keep a balanced fallback (square-ish).
    aspect: float | None = None
    if img_path is not None:
        try:
            with Image.open(img_path) as im:
                w_px, h_px = im.size
            if w_px > 0 and h_px > 0:
                aspect = w_px / h_px
        except (UnidentifiedImageError, OSError) as exc:
            log.warning("could not read image dimensions for %s: %s", file_id, exc)

    if aspect is None:
        # Placeholder / unknown — pick a 4:3 box (golden-friendly default).
        aspect = 4.0 / 3.0

    text_w_in, img_w_in, img_h_in, img_top_off_in = _image_box_for_aspect(aspect)

    # Bullets on the left, width derived from layout.
    _draw_bullets(
        slide, theme, section,
        left=Inches(_CONTENT_LEFT),
        top=Inches(_CONTENT_TOP),
        width=Inches(text_w_in),
        height=Inches(_CONTENT_HEIGHT),
    )

    img_left = Inches(_CONTENT_LEFT + text_w_in + _GUTTER)
    img_top = Inches(_CONTENT_TOP + img_top_off_in)
    img_w = Inches(img_w_in)
    img_h = Inches(img_h_in)

    if img_path is not None:
        try:
            # Pass width AND height computed from the real aspect — python-pptx
            # will honour them exactly, so no stretching occurs.
            slide.shapes.add_picture(str(img_path), img_left, img_top, width=img_w, height=img_h)
        except Exception as exc:  # noqa: BLE001
            log.warning("add_picture failed for %s: %s", file_id, exc)
            img_path = None

    if img_path is None:
        _add_rect(slide, img_left, img_top, img_w, img_h, theme.text_mute)
        label = section.image.caption if (section.image and section.image.caption) else (
            f"[image: {file_id}]" if file_id else "[image placeholder]"
        )
        _add_text(
            slide, img_left, img_top + Inches(img_h_in / 2 - 0.3), img_w, Inches(0.6),
            text=label,
            font=theme.body_font, size_pt=14, color=theme.neutral,
            align=PP_ALIGN.CENTER,
        )

    if section.image and section.image.caption and img_path is not None:
        _add_text(
            slide, img_left, img_top + img_h + Inches(0.05), img_w, Inches(0.4),
            text=section.image.caption,
            font=theme.body_font, size_pt=12, color=theme.text_mute,
            align=PP_ALIGN.CENTER,
        )


# Slide width in inches (16:9). Used for centring tables.
_SLIDE_WIDTH = 13.333


def _table_column_widths(
    headers: list[str], rows: list[list[str]], *, max_total_in: float = 11.9,
) -> list[float]:
    """Compute per-column widths in inches, sized to content.

    Heuristic: each column gets ``longest_cell_chars × 0.11"`` plus a small
    padding, clamped to ``[0.9", 4.5"]``. If the sum exceeds ``max_total_in``
    we scale all columns down proportionally so the table still fits the slide.
    """
    n = len(headers)
    raw: list[float] = []
    for ci in range(n):
        longest = len(headers[ci]) if ci < len(headers) else 1
        for r in rows:
            if ci < len(r):
                longest = max(longest, len(r[ci]))
        # Char-to-inch heuristic for ~12pt body / 14pt bold header.
        desired = longest * 0.11 + 0.4
        raw.append(max(0.9, min(4.5, desired)))
    total = sum(raw) or 1.0
    if total > max_total_in:
        scale = max_total_in / total
        raw = [w * scale for w in raw]
    return raw


def _table_section_slide(slide, theme: _Theme, section: Section) -> None:
    table = section.table
    if table is None or not table.headers:
        # Caller checks layout, but guard anyway — degrade to bullets.
        _bullets_section_slide(slide, theme, section)
        return

    n_cols = len(table.headers)
    # Cap rows at what fits visually; users can split into multiple sections
    # for huge tables. 14 rows × 28pt + header ≈ 4.4 inches.
    max_rows = 14
    rows = table.rows[:max_rows]
    n_rows = len(rows) + 1  # +1 for header

    col_widths_in = _table_column_widths(table.headers, rows)
    total_w_in = sum(col_widths_in)
    # Centre the table horizontally on the slide for visual balance.
    left_in = max(0.7, (_SLIDE_WIDTH - total_w_in) / 2.0)
    left, top = Inches(left_in), Inches(2.45)
    width = Inches(total_w_in)
    # Cap height so very tall tables don't run off the slide.
    height = Inches(min(0.45 * n_rows, 4.4))

    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = shape.table

    # Apply per-column widths so columns adapt to content rather than being
    # uniformly stretched.
    for ci, w_in in enumerate(col_widths_in):
        tbl.columns[ci].width = Inches(w_in)

    for ci, header in enumerate(table.headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = theme.primary
        tf = cell.text_frame
        tf.text = header
        for p in tf.paragraphs:
            for r in p.runs:
                r.font.name = theme.heading_font
                r.font.size = Pt(14)
                r.font.bold = True
                r.font.color.rgb = theme.text

    for ri, row in enumerate(rows, start=1):
        for ci in range(n_cols):
            cell = tbl.cell(ri, ci)
            tf = cell.text_frame
            tf.text = row[ci] if ci < len(row) else ""
            for p in tf.paragraphs:
                for r in p.runs:
                    r.font.name = theme.body_font
                    r.font.size = Pt(12)
                    r.font.color.rgb = theme.primary

    if table.caption:
        _add_text(
            slide, left, top + height + Inches(0.1), width, Inches(0.4),
            text=table.caption,
            font=theme.body_font, size_pt=12, color=theme.text_mute,
        )

    if len(table.rows) > max_rows:
        _add_text(
            slide, left, top + height + Inches(0.5), width, Inches(0.35),
            text=f"… {len(table.rows) - max_rows} more rows truncated",
            font=theme.body_font, size_pt=10, color=theme.text_mute,
        )


def _pick_layout(section: Section) -> str:
    """Return one of ``content-bullets`` / ``content-image`` / ``content-table``.

    Explicit ``layout_hint`` wins; else we infer from which payload is set.
    Mismatched hints (e.g. ``content-table`` without ``table``) gracefully fall
    back to bullets — the renderer should never raise on shape of input.
    """
    if section.layout_hint == "content-image" and section.image is not None:
        return "content-image"
    if section.layout_hint == "content-table" and section.table is not None:
        return "content-table"
    if section.layout_hint == "content-bullets":
        return "content-bullets"
    # Auto-infer
    if section.image is not None:
        return "content-image"
    if section.table is not None:
        return "content-table"
    return "content-bullets"


def _section_slide(
    prs: Presentation, theme: _Theme, idx: int, total: int, section: Section,
    uploads_dir: Path | None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, theme.neutral)
    _section_header(slide, theme, idx, total, section.heading)

    layout = _pick_layout(section)
    if layout == "content-image":
        _image_section_slide(slide, theme, section, uploads_dir)
    elif layout == "content-table":
        _table_section_slide(slide, theme, section)
    else:
        _bullets_section_slide(slide, theme, section)

    # Speaker notes go to the slide notes pane (PowerPoint reads it on the
    # presenter screen). Doesn't print, doesn't pollute the visible page.
    if section.speaker_notes:
        notes_tf = slide.notes_slide.notes_text_frame
        notes_tf.text = section.speaker_notes


def _closing_slide(prs: Presentation, theme: _Theme, doc: OutlineDoc) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, theme.primary)
    _add_rect(slide, Inches(1.0), Inches(2.4), Inches(0.6), Inches(0.06), theme.accent)
    label = "Thank you" if doc.language == "en" else "谢谢观看"
    _add_text(
        slide, Inches(1.0), Inches(2.55), Inches(11.3), Inches(1.4),
        text=label,
        font=theme.heading_font, size_pt=54, color=theme.text, bold=True,
    )
    _add_text(
        slide, Inches(1.0), Inches(4.0), Inches(11.3), Inches(0.8),
        text=doc.title,
        font=theme.body_font, size_pt=22, color=theme.text_mute,
    )


def render_outline(
    doc: OutlineDoc,
    template: TemplateInfo,
    *,
    uploads_dir: Path | None = None,
) -> bytes:
    """Build a .pptx for ``doc`` styled by ``template``; return the file bytes.

    ``uploads_dir`` is the directory where ``ImageRef.file_id`` resolves to a
    file on disk. ``None`` (the default) makes every image render as a
    placeholder rectangle — handy for tests and outline-only previews.
    """
    theme = _Theme.from_template(template)

    prs = Presentation()
    prs.slide_width = _SLIDE_W
    prs.slide_height = _SLIDE_H

    _cover_slide(prs, theme, doc)
    total = len(doc.sections)
    for i, section in enumerate(doc.sections, start=1):
        _section_slide(prs, theme, i, total, section, uploads_dir)
    _closing_slide(prs, theme, doc)

    buf = io.BytesIO()
    prs.save(buf)
    log.info(
        "render done: template=%s sections=%d total_slides=%d size_kb=%d",
        template.name, total, total + 2, len(buf.getvalue()) // 1024,
    )
    return buf.getvalue()
